import os
import time
import re
import pytz

from urllib.request import Request, urlopen
from bs4 import BeautifulSoup
from google.cloud import storage
import google.auth
from google.auth import transport
from datetime import datetime, timedelta
from data.db_connection import graph
from typing import Optional
import pdfplumber
import cfscrape
from pptx import Presentation
from github import Github


# Set properties for file uploads
uploadpath = './static/uploads' # LOCAL: Save the file to ./uploads
UPLOAD_FOLDER = uploadpath # LOCAL: Save the file to ./uploads

GITHUB_TOKEN = os.getenv('GITHUB_TOKEN',False)


def save_file(file,file_name,datapath,cloud=False,bucket_name=None):
    if (cloud==True) or (cloud=="True"):
        # https://cloud.google.com/appengine/docs/flexible/python/using-cloud-storage
        print('saving to cloud')
        blob_name = os.path.join(datapath,file_name)
        # file_path = os.path.join(datapath,blob_name)
        file_path = blob_name
        # upload_file_to_bucket(bucket_name, blob_name, file_path)
        bucket = storage.Client().bucket(bucket_name)
        blob = bucket.blob(blob_name)
        try:
            blob.upload_from_string(file.read(),content_type=file.content_type)
        except Exception as e:
            print(e)
    else:
        datapath = os.path.join(UPLOAD_FOLDER,datapath)
        if not os.path.exists(datapath):
            os.makedirs(datapath)

        file_path = os.path.join(datapath,file_name)
        print(file_path)
        file.save(file_path)

    return file_path



def getTextFromFile(file,url=False):
    filetext = ''
    if url:
        with pdfplumber.open(file) as pdf:
            for i,page in enumerate(pdf.pages):
                filetext += page.extract_text(x_tolerance=1,layout=False)

    else:
        if file.filename.endswith('.pdf'):
            with pdfplumber.open(file) as pdf:
                for i,page in enumerate(pdf.pages):
                    filetext += page.extract_text(x_tolerance=1,layout=False)

        elif file.filename.endswith('pptx') or file.filename.endswith('ppt'):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        addtext = ' ' + shape.text
                        filetext += addtext

    return filetext

def get_text_url(url):
    # Scrape
    for attempt in range(5):
        try:
            scraper = cfscrape.create_scraper()
            article = scraper.get(url).content
        except:
            time.sleep(2)
            continue
        else:
            break
    if article is None:
        article = ""
        title = None
    # Get body text and title
    soup = BeautifulSoup(article, features="html.parser")

    # Extract body text
    body = soup.find('body')
    findclasses = ['p','h1','h2','h3','h4','h5','h6','li','pre']
    bodytext = list(body.find_all(findclasses))
    # Remove everything before first title
    startelement = 0
    for i,tag in enumerate(bodytext):
        if tag.name in ['h1','h2','h3']:
            startelement = i
            break
    bodytext = bodytext[startelement:]

    cleantext = []
    rawtext = []
    for tag in bodytext:
        text = tag.text.strip()
        if len(text) > 10:
            htmltext = str(tag).split('>')[1:]
            htmltext = '>'.join(i for i in htmltext)
            htmltext = str(htmltext).split('<')[:-1]
            htmltext = '<'.join(i for i in htmltext)
            formatted_tag = f'<{tag.name}>{htmltext}</{tag.name}>'
            cleantext.append(formatted_tag)
            if tag.name == 'p':
                cleantext.append('<br>')
            rawtext.append(text)
    cleantext = ''.join(cleantext)
    rawtext = ' '.join(rawtext)

    return rawtext

def getTextFromGithub(url):
    print('getting from github')
    # Get text from github
    try:
        g = Github(GITHUB_TOKEN)
        # extract repo from url link
        if ('/blob/' in url):
            reponame = url.split('/blob/')[0].split('.com/')[-1]
        elif ('/tree/' in url):
            reponame = url.split('/tree/')[0].split('.com/')[-1]
        print(reponame)

        # Get the file path
        if 'main' in url:
            filepath = url.split('main/')[1]
        else:
            filepath = url.split('master/')[1]

        # Extract contents
        contents = g.get_repo(reponame).get_contents(filepath)
        repopath = reponame + '/' + filepath

        filenames = []
        codes = []
        languages = []

        # If contents is a directory, get all code files in the directory
        if type(contents) == list:
            print('This is a directory')
            for file in contents:
                print(file.path)
                if file.path.endswith('.py'):
                    filename = file.name
                    code = file.decoded_content.decode()
                    code = code.strip()
                    language = 'python'
                    filenames.append(filename)
                    codes.append(code)
                    languages.append(language)

                elif file.path.lower().endswith('.r') or file.path.lower().endswith('.rmd'):
                    filename = file.name
                    # print(file.decoded_content.decode())
                    code = file.decoded_content.decode()
                    code = code.strip()
                    language = 'R'
                    filenames.append(filename)
                    codes.append(code)
                    languages.append(language)

                else:
                    ending = file.name.split('.')[-1]
                    print('ending')
                    print(ending)
                    if ending in ['txt','md','yaml','Makefile','Dockerfile']:
                        filename = file.name
                        code = file.decoded_content.decode()
                        code = code.strip()
                        language = 'other'
                        filenames.append(filename)
                        codes.append(code)
                        languages.append(language)
                 
        else:
            if contents.path.endswith('.py'):
                filename = contents.name
                code = contents.decoded_content.decode()
                # strip whitespace in front and back of code
                code = code.strip()
                language = 'python'
                filenames.append(filename)
                codes.append(code)
                languages.append(language)
            elif contents.path.lower().endswith('.r') or contents.path.lower().endswith('.rmd'):
                filename = contents.name
                code = contents.decoded_content.decode()
                code = code.strip()
                language = 'R'
                filenames.append(filename)
                codes.append(code)
                languages.append(language)
            

        return repopath,filenames,codes,languages

    except Exception as e:
        print(e)
        return [],[],[],[]
    

def get_signed_file_link(file_name,bucket_name=None):
    # https://stackoverflow.com/questions/64234214/how-to-generate-a-blob-signed-url-in-google-cloud-run/64245028#64245028
    # Set CORS on bucket: https://cloud.google.com/storage/docs/gsutil/commands/cors#synopsis
    
    credentials, project_id = google.auth.default()

    # Perform a refresh request to get the access token of the current credentials (Else, it's None)
    r = transport.requests.Request()
    credentials.refresh(r)

    client = storage.Client()
    bucket = client.bucket(bucket_name)
    blob_name = file_name
    blob = bucket.blob(blob_name)
    expires = datetime.now() + timedelta(seconds=3600)

    # In case of user credential use, define manually the service account to use (for development purpose only)
    service_account_email = "YOUR DEV SERVICE ACCOUNT"
    # If you use a service account credential, you can use the embedded email
    if hasattr(credentials, "service_account_email"):
        service_account_email = credentials.service_account_email

    url = blob.generate_signed_url(expiration=expires,service_account_email=service_account_email, access_token=credentials.token)

    return url

def convertDeadlineToUTC(deadline,timezone):
    if len(deadline) > 0:
        deadline = datetime.strptime(deadline, '%m/%d/%Y %H:%M')
        if timezone is not None:
            deadline = pytz.timezone(timezone).localize(deadline, is_dst=None)
        else:
            deadline = pytz.timezone('US/Eastern').localize(deadline, is_dst=None)
        deadline = deadline.astimezone(pytz.utc)
        deadline = str(deadline.strftime('%m/%d/%Y %H:%M'))

    return deadline

def convertDeadlineToLocal(deadline,timezone):
    if len(deadline) > 0:
        deadline = datetime.strptime(deadline, '%m/%d/%Y %H:%M')
    #    localize timezone as UTC
        deadline = pytz.timezone('UTC').localize(deadline, is_dst=None)
    #    convert to local timezone
        if timezone is not None:
            deadline = deadline.astimezone(pytz.timezone(timezone))
        else:
            deadline = deadline.astimezone(pytz.timezone('US/Eastern'))

        deadline = str(deadline.strftime('%m/%d %H:%M'))

    return deadline

def convertTimeToBulma(deadline,timezone):
    if len(deadline) > 0:
        deadline = datetime.strptime(deadline, '%m/%d/%Y %H:%M')
    #    localize timezone as UTC
        deadline = pytz.timezone('UTC').localize(deadline, is_dst=None)
    #    convert to local timezone
        if timezone is not None:
            deadline = deadline.astimezone(pytz.timezone(timezone))
        else:
            deadline = deadline.astimezone(pytz.timezone('US/Eastern'))

        deadline = str(deadline.strftime('%Y-%m-%d %H:%M'))

    return deadline